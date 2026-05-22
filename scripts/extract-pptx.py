"""
Extract all text, images, and tables from the Rewari inventory PPTX.
Produces:
  - public/extracted-photos/slide_<N>_img_<M>.<ext>
  - docs/extracted-listings.json
"""

import json, re, os
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH  = Path("docs/REWARI 24.03.2026 (1).pptx")
PHOTOS_DIR = Path("public/extracted-photos")
OUT_JSON   = Path("docs/extracted-listings.json")

PHOTOS_DIR.mkdir(parents=True, exist_ok=True)

# ── text helpers ──────────────────────────────────────────────────────────────

def shape_text(shape):
    try:
        if shape.has_text_frame:
            lines = []
            for para in shape.text_frame.paragraphs:
                t = "".join(r.text for r in para.runs).strip()
                if t:
                    lines.append(t)
            return lines
    except Exception:
        pass
    return []

def table_to_rows(table):
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        if any(cells):
            rows.append(cells)
    return rows

# ── title parsing (all titles follow: CITY-LOCATION SIZE [NL]) ────────────────

def parse_title(title):
    """
    Parse agency PPTX slide title into structured fields.
    Pattern: CITY-DESCRIPTION WIDTHXHEIGHT [NL]
    Examples:
      REWARI-RAILWAY STATION PLATFORM 20x10 NL
      HUDA BYPASS CHOWK 20x10
      OPP BESTECH FLATS 20x10 NL
    """
    t = title.strip()

    # illumination: NL = non-lit
    illuminated = True
    illu_type   = "frontlit"  # default assumption if not NL
    if re.search(r'\bNL\b', t, re.I):
        illuminated = False
        illu_type   = ""
        t = re.sub(r'\bNL\b', '', t, flags=re.I).strip()

    # size
    m = re.search(r'(\d{2,3})\s*[xX]\s*(\d{2,3})', t)
    w = int(m.group(1)) if m else 0
    h = int(m.group(2)) if m else 0
    if m:
        t = t[:m.start()].strip()  # remove size from remaining text

    # strip leading city prefix e.g. "REWARI-", "DHARUHERA-"
    city = "Rewari"  # default; overridden by section detection
    t = re.sub(r'^(REWARI|DHARUHERA)\s*[-–]\s*', '', t, flags=re.I).strip()

    # landmark: NEAR X, OPP X, NR X
    landmark = ""
    m2 = re.search(r'\b(?:NEAR|NR|OPP|OPPOSITE)\s+(.+)$', t, re.I)
    if m2:
        landmark = m2.group(1).strip().title()
        t = t[:m2.start()].strip()

    location_clean = t.title()  # e.g. "Railway Station Platform"

    return {
        "location_clean": location_clean,
        "landmark":       landmark,
        "size_w":         w,
        "size_h":         h,
        "illuminated":    illuminated,
        "illu_type":      illu_type,
    }

# ── locality & type mapping ───────────────────────────────────────────────────

# Map location keywords → locality slug
LOCALITY_MAP = [
    # Known localities in site's 15-slug list
    (["railway station", "nai basti railway"],  "rewari-railway-station"),
    (["delhi road"],                             "delhi-road"),
    (["sadar bazar", "sadar market"],            "sadar-bazar"),
    (["bus stand", "hrtc"],                      "rewari-bus-stand"),
    (["gurudwara chowk"],                        "gurudwara-chowk"),
    (["grain market", "anaj mandi"],             "new-grain-market"),
    (["majra chowk", "majra"],                   "majra-chowk"),
    (["jat college"],                            "jat-college-road"),
    (["circular road"],                          "circular-road"),
    (["model town"],                             "model-town"),
    (["subhash nagar", "subash chowk"],          "subhash-nagar"),
    (["industrial area", "imt"],                 "rewari-industrial-area"),
    (["kosli"],                                  "kosli-road"),
    (["nh-48", "nh48", "nh 48", "bypass", "huda bypass"], "nh-48"),
    # Partial/approximate mappings — flag for manual review
    (["mahendergarh road"],                      "dharan-road"),   # south exit, closest match
    (["bawal road"],                             "rewari-industrial-area"),  # Bawal direction/industrial belt
    (["piolet", "pilot chowk"],                  "majra-chowk"),   # inner city chowk, closest match
    (["jhajjar road"],                           "kosli-road"),    # south exit corridor, flag for review
    # Locations with no close match — left blank for manual assignment
    (["pallawas", "kulana", "pataudi",
      "rampura", "bitwana"],                     ""),
]

def map_locality(location_clean, landmark, title_original, city):
    if city == "Dharuhera":
        return ""   # Dharuhera not in Rewari locality list
    combined = (location_clean + " " + landmark + " " + title_original).lower()
    for keywords, slug in LOCALITY_MAP:
        if any(kw in combined for kw in keywords):
            return slug
    return ""

# Type classification
def classify_type(location_clean, landmark, title_original):
    combined = (location_clean + " " + landmark + " " + title_original).lower()
    if "platform" in combined or "car parking" in combined or "parking" in combined:
        return "railway-station"
    if "railway" in combined or "station" in combined:
        return "railway-station"
    if "unipole" in combined:
        return "unipole"
    if "gantry" in combined:
        return "gantry"
    if "led" in combined:
        return "led-display"
    # All others are standard hoardings
    return "hoarding"

# Slug generation for listing
_slug_used = {}
def make_slug(title, city="rewari"):
    base = re.sub(r'[^a-z0-9 ]', '', title.lower())
    base = re.sub(r'\s+', '-', base.strip())
    base = f"{city.lower()}-{base}"
    base = base[:70]
    n = _slug_used.get(base, 0)
    _slug_used[base] = n + 1
    return base if n == 0 else f"{base}-{n}"

# ── raw extraction ─────────────────────────────────────────────────────────────

prs = Presentation(str(PPTX_PATH))
total_slides = len(prs.slides)

slide_data = []
for slide_idx, slide in enumerate(prs.slides, start=1):
    texts        = []
    tables       = []
    images_saved = []

    for shape_idx, shape in enumerate(slide.shapes, start=1):
        texts.extend(shape_text(shape))
        if shape.has_table:
            tables.append(table_to_rows(shape.table))
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img   = shape.image
                fname = f"slide_{slide_idx:02d}_img_{shape_idx:02d}.{img.ext}"
                (PHOTOS_DIR / fname).write_bytes(img.blob)
                images_saved.append(fname)
        except Exception:
            pass

    for ph in slide.placeholders:
        try:
            if hasattr(ph, "image"):
                img   = ph.image
                idx_n = ph.placeholder_format.idx
                fname = f"slide_{slide_idx:02d}_ph_{idx_n}.{img.ext}"
                if fname not in images_saved:
                    (PHOTOS_DIR / fname).write_bytes(img.blob)
                    images_saved.append(fname)
        except Exception:
            pass

    slide_data.append({
        "slide":  slide_idx,
        "texts":  texts,
        "tables": tables,
        "images": images_saved,
    })

# ── build structured listings ─────────────────────────────────────────────────

# Detect section dividers (slides with a single ALL-CAPS city word, no size)
SECTION_DIVIDERS = {}   # slide_idx -> city name
for sd in slide_data:
    if len(sd["texts"]) == 1 and not sd["images"]:
        t = sd["texts"][0].strip().upper()
        if t in ("REWARI", "DHARUHERA", "BAWAL", "NARNAUL"):
            SECTION_DIVIDERS[sd["slide"]] = t.title()

listings = []
flags    = []
current_city = "Rewari"

for sd in slide_data:
    idx   = sd["slide"]
    texts = sd["texts"]

    # Update city from section divider
    if idx in SECTION_DIVIDERS:
        current_city = SECTION_DIVIDERS[idx]
        flags.append({"slide": idx, "reason": f"Section divider — city set to '{current_city}'"})
        continue

    if not texts:
        flags.append({"slide": idx, "reason": "No text found — image-only slide, skipping"})
        continue

    # Primary title = first text item
    raw_title = texts[0]
    parsed    = parse_title(raw_title)

    loc_clean = parsed["location_clean"]
    landmark  = parsed["landmark"]
    w         = parsed["size_w"]
    h         = parsed["size_h"]

    otype    = classify_type(loc_clean, landmark, raw_title)
    locality = map_locality(loc_clean, landmark, raw_title, current_city)

    # Build a human-readable title
    if landmark:
        display_title = f"{loc_clean} — {landmark}, {current_city}"
    else:
        display_title = f"{loc_clean}, {current_city}"

    listing = {
        "title":             display_title,
        "type":              otype,
        "locality":          locality,
        "city":              current_city,
        "landmark":          landmark,
        "size_width_ft":     w,
        "size_height_ft":    h,
        "facing":            "",   # not in source data
        "illuminated":       parsed["illuminated"],
        "illumination_type": parsed["illu_type"],
        "structure_type":    "",   # not in source data
        "traffic_count":     "",   # not in source data
        "availability":      "available",
        "is_featured":       False,
        "notes_internal":    f"Extracted from agency PPTX — slide {idx} — {raw_title}",
        "_slide":            idx,
        "_raw_title":        raw_title,
        "_images":           sd["images"],
        "_extra_texts":      texts[1:],
    }
    listings.append(listing)

    # Flags for missing key fields
    missing = []
    if not listing["type"]:     missing.append("type")
    if not listing["locality"]: missing.append("locality")
    if not listing["size_width_ft"]: missing.append("size")
    if missing:
        flags.append({
            "slide":   idx,
            "title":   raw_title,
            "city":    current_city,
            "missing": missing,
        })

# ── write JSON ────────────────────────────────────────────────────────────────

output = {
    "meta": {
        "source_file":       str(PPTX_PATH),
        "total_slides":      total_slides,
        "section_dividers":  list(SECTION_DIVIDERS.values()),
        "listings_extracted": len(listings),
        "images_extracted":  sum(len(l["_images"]) for l in listings),
    },
    "listings": listings,
    "flags":    flags,
}
OUT_JSON.write_text(json.dumps(output, indent=2, ensure_ascii=False), encoding="utf-8")

# ── print summary ─────────────────────────────────────────────────────────────

print(f"Total slides:       {total_slides}")
print(f"Section dividers:   {list(SECTION_DIVIDERS.items())}")
print(f"Listings extracted: {len(listings)}")
print(f"Images extracted:   {output['meta']['images_extracted']}")
print()

fields = ["type","locality","city","size_width_ft","illuminated","landmark"]
print("Field population rate:")
for f in fields:
    filled = sum(1 for l in listings if l.get(f) not in (None,"",0,False))
    pct    = filled/len(listings)*100 if listings else 0
    print(f"  {f:<22} {filled:>3}/{len(listings)}  ({pct:.0f}%)")

print()
print(f"Slides flagged for review: {len(flags)}")
for flag in flags:
    s       = flag['slide']
    reason  = flag.get('reason') or f"missing: {flag.get('missing', [])}"
    title   = flag.get('title','')[:50]
    city    = flag.get('city','')
    print(f"  Slide {s:>2} [{city:<10}]: {title:<50}  [{reason}]")

print()
print("All listing titles:")
for l in listings:
    print(f"  [{l['_slide']:>2}] {l['city']:<10} {l['type']:<20} {l['size_width_ft']}x{l['size_height_ft']}  NL={not l['illuminated']}  {l['_raw_title']}")
